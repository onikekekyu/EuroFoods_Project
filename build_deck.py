#!/usr/bin/env python3
"""EuroFoods S.A. — FP&A presentation deck (Lovable palette, white)."""
import json
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION

D = json.load(open("data/dashboard_data.json"))

# ---- palette ----
INK=RGBColor(0x1A,0x16,0x25); MUT=RGBColor(0x6B,0x64,0x78); LINE=RGBColor(0xEC,0xE9,0xE4)
BG=RGBColor(0xFF,0xFF,0xFF); PANEL=RGBColor(0xFB,0xFA,0xF9)
CYAN=RGBColor(0x06,0xB6,0xD4); PINK=RGBColor(0xFF,0x4D,0x8D); ORANGE=RGBColor(0xFF,0x7A,0x45)
AMBER=RGBColor(0xFF,0x9E,0x2C); VIOLET=RGBColor(0x7C,0x5C,0xFF); GREEN=RGBColor(0x16,0xB3,0x64)
RED=RGBColor(0xF0,0x38,0x4B); GREEND=RGBColor(0x0B,0x7A,0x45); CYAND=RGBColor(0x0E,0x74,0x90)
WHITE=RGBColor(0xFF,0xFF,0xFF)
FONT="Inter"; FONTF=("Inter","Segoe UI","Calibri")

prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
BL=prs.slide_layouts[6]
SW,SH=13.333,7.5

def slide():
    s=prs.slides.add_slide(BL)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=BG; r.line.fill.background()
    r.shadow.inherit=False
    return s

def box(s,x,y,w,h,fill=None,line=None,lw=0.75,radius=True,shadow=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                           I(x),I(y),I(w),I(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    if radius:
        try: shp.adjustments[0]=0.06
        except: pass
    return shp

def grad(shp,c1=PINK,c2=ORANGE,angle=0):
    shp.fill.gradient()
    try:
        shp.fill.gradient_angle=angle
        gs=shp.fill.gradient_stops
        gs[0].position=0.0; gs[0].color.rgb=c1
        gs[1].position=1.0; gs[1].color.rgb=c2
    except Exception:
        shp.fill.solid(); shp.fill.fore_color.rgb=c1
    shp.line.fill.background(); shp.shadow.inherit=False
    return shp

def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.0):
    tb=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=sp
        for (t,sz,col,bold) in para:
            r=p.add_run(); r.text=t; f=r.font
            f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.name=FONT
    return tb

def R(t,sz,col,bold=False): return (t,sz,col,bold)

def logo(s,x,y,sz=0.62):
    g=box(s,x,y,sz,sz,fill=PINK,radius=True); grad(g,PINK,ORANGE,45)
    try: g.adjustments[0]=0.28
    except: pass
    txt(s,x,y,sz,sz,[[R("€F",15,WHITE,True)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)

def header(s,kicker,title,n):
    grad(box(s,0,0,0.16,SH,fill=PINK),PINK,ORANGE,90)
    logo(s,0.55,0.5,0.6)
    txt(s,1.32,0.5,9,0.3,[[R(kicker,11,CYAND,True)]])
    txt(s,1.32,0.78,10.5,0.6,[[R(title,25,INK,True)]])
    txt(s,12.0,0.55,0.9,0.4,[[R(f"{n:02d}",13,RGBColor(0xC9,0xC4,0xCF),True)]],PP_ALIGN.RIGHT)
    box(s,1.32,1.42,10.9,0.018,fill=LINE,radius=False)

def kpi(s,x,y,w,h,lab,val,sub,accent,subcol=None):
    box(s,x,y,w,h,fill=PANEL,line=LINE,lw=0.75)
    box(s,x,y,0.07,h,fill=accent,radius=False)
    txt(s,x+0.22,y+0.16,w-0.3,0.3,[[R(lab.upper(),9.5,MUT,True)]])
    txt(s,x+0.22,y+0.42,w-0.3,0.55,[[R(val,27,INK,True)]])
    txt(s,x+0.22,y+h-0.42,w-0.3,0.3,[[R(sub,10,subcol or MUT,bool(subcol))]])

def chart_style(ch,sz=10):
    ch.has_title=False
    try:
        ch.font.size=Pt(sz); ch.font.name=FONT; ch.font.color.rgb=MUT
    except: pass

def add_bar(s,x,y,w,h,cats,series,colors,legend=True,stacked=False,horiz=False,fmt='#,##0'):
    cd=CategoryChartData(); cd.categories=cats
    for nm,vals in series: cd.add_series(nm,vals)
    ctype=XL_CHART_TYPE.BAR_CLUSTERED if horiz else XL_CHART_TYPE.COLUMN_CLUSTERED
    if stacked: ctype=XL_CHART_TYPE.BAR_STACKED if horiz else XL_CHART_TYPE.COLUMN_STACKED
    gf=s.shapes.add_chart(ctype,I(x),I(y),I(w),I(h),cd); ch=gf.chart
    chart_style(ch)
    if legend and len(series)>1:
        ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.TOP; ch.legend.include_in_layout=False
        ch.legend.font.size=Pt(10); ch.legend.font.color.rgb=INK
    else: ch.has_legend=False
    for i,pl in enumerate(ch.plots[0].series):
        if len(series)==1:
            for j,pt in enumerate(pl.points):
                pt.format.fill.solid(); pt.format.fill.fore_color.rgb=colors[j%len(colors)]
        else:
            pl.format.fill.solid(); pl.format.fill.fore_color.rgb=colors[i%len(colors)]
        pl.format.line.fill.background()
    try:
        ch.value_axis.has_major_gridlines=True
        ch.value_axis.major_gridlines.format.line.color.rgb=LINE
        ch.value_axis.major_gridlines.format.line.width=Pt(0.5)
        ch.value_axis.format.line.fill.background()
        ch.value_axis.tick_labels.font.size=Pt(9); ch.value_axis.tick_labels.font.color.rgb=MUT
        ch.value_axis.tick_labels.number_format=fmt; ch.value_axis.tick_labels.number_format_is_linked=False
        ch.category_axis.tick_labels.font.size=Pt(10); ch.category_axis.tick_labels.font.color.rgb=INK
        ch.category_axis.format.line.color.rgb=LINE
    except: pass
    ch.plots[0].gap_width=80
    return ch

def add_line(s,x,y,w,h,cats,series,colors,fmt='#,##0'):
    cd=CategoryChartData(); cd.categories=cats
    for nm,vals in series: cd.add_series(nm,vals)
    gf=s.shapes.add_chart(XL_CHART_TYPE.LINE,I(x),I(y),I(w),I(h),cd); ch=gf.chart
    chart_style(ch)
    if len(series)>1:
        ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.TOP; ch.legend.include_in_layout=False
        ch.legend.font.size=Pt(10); ch.legend.font.color.rgb=INK
    else: ch.has_legend=False
    for i,pl in enumerate(ch.series):
        pl.format.line.color.rgb=colors[i%len(colors)]; pl.format.line.width=Pt(2.5)
        pl.smooth=True
    try:
        ch.value_axis.has_major_gridlines=True
        ch.value_axis.major_gridlines.format.line.color.rgb=LINE
        ch.value_axis.major_gridlines.format.line.width=Pt(0.5)
        ch.value_axis.format.line.fill.background()
        ch.value_axis.tick_labels.font.size=Pt(9); ch.value_axis.tick_labels.font.color.rgb=MUT
        ch.value_axis.tick_labels.number_format=fmt; ch.value_axis.tick_labels.number_format_is_linked=False
        ch.category_axis.tick_labels.font.size=Pt(8.5); ch.category_axis.tick_labels.font.color.rgb=MUT
        ch.category_axis.format.line.color.rgb=LINE
    except: pass
    return ch

def chip(s,x,y,txt_,col,w=1.5):
    c=box(s,x,y,w,0.34,fill=col,radius=True)
    try: c.adjustments[0]=0.5
    except: pass
    txt(s,x,y+0.02,w,0.3,[[R(txt_,9.5,WHITE,True)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)

EUR=lambda n:('€%.2fM'%(n/1e6)) if abs(n)>=1e6 else ('€%.0fk'%(n/1e3))

# =================================================================== SLIDE 1 — TITLE
s=slide()
grad(box(s,0,0,SW,SH,fill=PINK),RGBColor(0xFF,0xFF,0xFF),RGBColor(0xFF,0xFF,0xFF))  # keep white
box(s,0,0,SW,SH,fill=BG,radius=False)
# soft accent blobs
grad(box(s,9.2,-1.6,5.5,5.5,fill=PINK),RGBColor(0xFF,0xE3,0xEC),RGBColor(0xFF,0xFF,0xFF),45)
grad(box(s,-1.8,4.4,5.5,5.5,fill=CYAN),RGBColor(0xE1,0xF7,0xFB),RGBColor(0xFF,0xFF,0xFF),45)
logo(s,1.0,0.95,1.0)
txt(s,2.25,1.06,8,0.4,[[R("EuroFoods S.A.",17,INK,True)]])
txt(s,2.25,1.5,8,0.3,[[R("DATA FOR FP&A · 2026",11,MUT,True)]])
txt(s,1.0,2.7,11.6,1.7,[
   [R("Healthy growth, ",40,INK,True),R("missed plan",40,PINK,True),R(".",40,INK,True)],
   [R("A two-year FP&A diagnosis of EuroFoods.",40,INK,True)],
],sp=1.02)
txt(s,1.02,4.95,11,0.5,[[R("+5.4% revenue growth · 38.7% margin held · but 4.6% below budget in every market",15,MUT,False)]])
# mini kpi row
for i,(l,v,c) in enumerate([("Revenue 2024","€2.36M",CYAN),("YoY Growth","+5.4%",GREEN),
        ("Budget Var.","−4.6%",RED),("Margin","38.7%",VIOLET)]):
    x=1.0+i*2.55
    box(s,x,5.7,2.3,0.95,fill=PANEL,line=LINE)
    box(s,x,5.7,0.06,0.95,fill=c,radius=False)
    txt(s,x+0.2,5.82,2.0,0.25,[[R(l.upper(),8.5,MUT,True)]])
    txt(s,x+0.2,6.08,2.0,0.45,[[R(v,20,INK,True)]])
txt(s,1.0,6.95,11,0.3,[[R("FP&A Analytics  ·  Power BI · DAX · Scenario modelling  ·  Data as of December 2024",10,MUT,False)]])

# =================================================================== SLIDE 2 — CONTEXT / APPROACH
s=slide(); header(s,"CONTEXT & APPROACH","Two years of data, one star schema",2)
txt(s,1.32,1.62,5.4,0.4,[[R("THE BUSINESS",11,CYAND,True)]])
facts=[("3 countries","France · Germany · UK"),("20 SKUs","across 4 categories"),
       ("4 channels","Hyper · Super · Conv · Online"),("36 customers","named retail accounts"),
       ("25,200 rows","weekly, Jan 2023 – Dec 2024")]
for i,(a,b) in enumerate(facts):
    y=2.0+i*0.78
    box(s,1.32,y,5.2,0.66,fill=PANEL,line=LINE)
    txt(s,1.54,y+0.07,2.4,0.5,[[R(a,15,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,3.7,y+0.07,2.7,0.5,[[R(b,11,MUT,False)]],anchor=MSO_ANCHOR.MIDDLE)
txt(s,7.0,1.62,5.4,0.4,[[R("THE DATA MODEL",11,CYAND,True)]])
box(s,7.0,2.0,5.3,3.05,fill=PANEL,line=LINE)
# star schema
cx,cy=9.65,3.5
f=box(s,cx-0.85,cy-0.45,1.7,0.9,fill=CYAN); txt(s,cx-0.85,cy-0.45,1.7,0.9,[[R("Transactions",11,WHITE,True)],[R("25,200 · fact",8.5,WHITE,False)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
dims=[("Dim_Date",8.05,2.35,VIOLET),("Dim_Products",11.05,2.35,PINK),
      ("Dim_Customers",8.05,4.35,AMBER),("Budget",11.05,4.35,GREEN)]
for nm,x,y,c in dims:
    box(s,x-0.7,y-0.32,1.5,0.64,fill=WHITE,line=c,lw=1.5)
    txt(s,x-0.7,y-0.32,1.5,0.64,[[R(nm,9.5,INK,True)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    ln=s.shapes.add_connector(2,I(x),I(y if y<cy else y),I(cx),I(cy)); ln.line.color.rgb=LINE; ln.line.width=Pt(1.25)
txt(s,7.0,5.2,5.3,0.5,[[R("Challenge solved: ",10,INK,True),R("transactions aggregated to (Year,Month,Country,Category) to join Budget's monthly grain.",10,MUT,False)]])
txt(s,1.32,5.95,11,1.0,[[R("DAX foundation:  ",11,CYAND,True),
   R("Total Revenue · Gross Margin · Margin % · Revenue PY (SAMEPERIODLASTYEAR) · YoY Growth · Revenue Variance · Forecast Revenue (What-If) — built on a calendar table marked as a date table.",11,MUT,False)]])

# =================================================================== SLIDE 3 — HEADLINE KPIs
s=slide(); header(s,"2024 PERFORMANCE","The headline: grew, profited, but under plan",3)
kdata=[("Revenue 2024","€2.36M","+5.4% vs 2023",CYAN,GREEND),
       ("Gross Margin","€915k","38.7% — held flat",GREEN,MUT),
       ("Budget Variance","−€115k","−4.6% unfavourable",RED,RED),
       ("Q3 2024 YoY","+11.9%","summer acceleration",VIOLET,GREEND)]
for i,(l,v,sub,acc,sc) in enumerate(kdata):
    kpi(s,1.32+i*2.78,1.75,2.62,1.5,l,v,sub,acc,sc)
# two-up charts
txt(s,1.32,3.5,5,0.3,[[R("Revenue by country — actual vs budget",12,INK,True)]])
cn=[c["Country"] for c in D["country"]]
act=[c["Rev"] for c in D["country"]]; bud=[c["Budget"] for c in D["country"]]
add_bar(s,1.32,3.85,5.5,3.0,cn,[("Actual",act),("Budget",bud)],[PINK,AMBER],fmt='€#,##0,,"M"')
txt(s,7.2,3.5,5,0.3,[[R("Monthly revenue — 2023 vs 2024",12,INK,True)]])
add_line(s,7.0,3.85,5.5,3.0,D["months"],[("2023",D["rev_by_month_23"]),("2024",D["rev_by_month_24"])],[RGBColor(0xC9,0xC4,0xCF),CYAN],fmt='€#,##0,"k"')

# =================================================================== SLIDE 4 — FINDING 1
s=slide(); header(s,"FINDING 01","The budget gap is volume, not price",4)
box(s,1.32,1.75,5.3,4.6,fill=PANEL,line=LINE)
txt(s,1.6,2.0,4.8,0.4,[[R("THE PRICE–VOLUME BRIDGE",10.5,CYAND,True)]])
txt(s,1.6,2.45,4.8,1.0,[[R("−€226k",46,RED,True)],[R("volume effect (selling fewer units)",11,MUT,False)]])
txt(s,1.6,3.75,4.8,1.0,[[R("≈ €0",40,INK,True)],[R("price / mix effect",11,MUT,False)]])
box(s,1.6,4.95,4.7,1.2,fill=WHITE,line=LINE)
txt(s,1.78,5.08,4.4,1.0,[[R("Actual prices landed on target. The entire shortfall is ~4.7% fewer units than planned — a planning error, not a discounting problem.",11,INK,False)]])
txt(s,7.0,1.75,5.3,0.3,[[R("Variance is uniform across every market (% vs budget)",11.5,INK,True)]])
add_bar(s,7.0,2.15,5.4,4.1,
   ["France","Germany","UK","Frozen","Dry","Fresh","Bev"],
   [("Var %",[-5.0,-5.0,-5.0,-5.0,-5.0,-4.7,-4.8])],
   [RED]*7,legend=False,fmt='0"%"')

# =================================================================== SLIDE 5 — FINDING 2
s=slide(); header(s,"FINDING 02","Online is the slowest channel, not the fastest",5)
seg=sorted(D["segment"],key=lambda x:-x["Growth%"])
names=[x["Segment"] for x in seg]; gr=[x["Growth%"] for x in seg]
cols=[GREEN if g>=5.4 else AMBER for g in gr]
txt(s,1.32,1.7,6,0.3,[[R("Revenue growth 2023 → 2024 by segment",12,INK,True)]])
add_bar(s,1.32,2.1,6.2,4.2,names,[("YoY %",gr)],cols,legend=False,fmt='0.0"%"')
box(s,8.0,1.9,4.3,4.3,fill=PANEL,line=LINE)
txt(s,8.25,2.15,3.85,0.4,[[R("WHY IT MATTERS",10.5,CYAND,True)]])
pts=[("Supermarket leads","+5.9% — 33.6% of revenue, the real growth engine"),
     ("Online lags","+4.9%, the slowest of all four channels"),
     ("Brick-and-mortar +6.0%","actually outgrew Online (+5.0%)"),
     ("The opportunity","Online is 22.6% of revenue but under-invested")]
for i,(a,b) in enumerate(pts):
    y=2.6+i*0.88
    box(s,8.25,y,0.12,0.66,fill=[GREEN,RED,CYAN,VIOLET][i],radius=False)
    txt(s,8.5,y,3.6,0.7,[[R(a+"  ",11,INK,True)],[R(b,9.5,MUT,False)]],sp=1.0)

# =================================================================== SLIDE 6 — FINDING 3 / MARGINS
s=slide(); header(s,"FINDING 03","Frozen Food is the structural margin drag",6)
txt(s,1.32,1.7,6,0.3,[[R("Margin % vs category target",12,INK,True)]])
cats=["Beverages","Fresh Produce","Dry Goods","Frozen Food"]
actm=[41.9,29.9,47.9,34.8]; tgt=[42,30,48,35]
add_bar(s,1.32,2.1,6.2,4.2,cats,[("Actual %",actm),("Target %",tgt)],[CYAN,RGBColor(0xC9,0xC4,0xCF)],fmt='0"%"')
box(s,8.0,1.9,4.3,4.3,fill=PANEL,line=LINE)
txt(s,8.25,2.15,3.85,0.4,[[R("THE FROZEN PROBLEM",10.5,CYAND,True)]])
txt(s,8.25,2.6,3.85,3.4,[
  [R("All 5 ",11,INK,True),R("Frozen Food SKUs sit below the 35% target.",11,MUT,False)],
  [R("",6,MUT,False)],
  [R("€1.42M",26,VIOLET,True)],
  [R("largest category by revenue — yet the lowest margin band.",10,MUT,False)],
  [R("",6,MUT,False)],
  [R("Verdict: ",11,INK,True),R("a category-wide sourcing issue, not weak individual products. No SKU warrants discontinuation.",10.5,MUT,False)],
],sp=1.05)

# =================================================================== SLIDE 7 — BUDGET VARIANCE
s=slide(); header(s,"BUDGET VARIANCE","Every market missed — by almost exactly the same %",7)
kpi(s,1.32,1.75,2.9,1.4,"Actual (2yr)","€4.60M","",CYAN)
kpi(s,4.42,1.75,2.9,1.4,"Budget (2yr)","€4.83M","",AMBER)
kpi(s,7.52,1.75,2.9,1.4,"Variance","−€227k","−4.7% unfavourable",RED,RED)
kpi(s,10.62,1.75,1.6,1.4,"Vol. var","−4.7%","",VIOLET)
txt(s,1.32,3.45,6,0.3,[[R("Actual vs budget revenue by country",12,INK,True)]])
cn=[c["Country"] for c in D["country"]]
add_bar(s,1.32,3.85,5.5,2.95,cn,[("Actual",[c["Rev"] for c in D["country"]]),("Budget",[c["Budget"] for c in D["country"]])],[PINK,AMBER],fmt='€#,##0,,"M"')
box(s,7.1,3.5,5.2,3.3,fill=PANEL,line=LINE)
txt(s,7.35,3.72,4.7,0.4,[[R("THE DIAGNOSIS",10.5,CYAND,True)]])
txt(s,7.35,4.15,4.75,2.5,[
  [R("Variance band: ",11,INK,True),R("−5.0% to −3.9% across all 3 countries × 4 categories × 8 quarters.",10.5,MUT,False)],
  [R("",6,MUT,False)],
  [R("Such uniformity is the signature of a ",10.5,MUT,False),R("top-down budget set ~5% too high",10.5,INK,True),R(" — not localised execution failures.",10.5,MUT,False)],
  [R("",6,MUT,False)],
  [R("Fix: ",11,GREEND,True),R("re-baseline 2025 on the real run-rate.",10.5,MUT,False)],
],sp=1.05)

# =================================================================== SLIDE 8 — SEASONALITY
s=slide(); header(s,"SEASONALITY","Beverages swing hardest — a planning signal",8)
si=D["seasonal_index"]
months=[r["Month"] for r in si]
mlab=["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
bev=[r["Beverages"] for r in si]; fro=[r["Frozen Food"] for r in si]
add_line(s,1.32,1.9,7.3,4.5,mlab,[("Beverages",bev),("Frozen Food",fro)],[CYAN,VIOLET],fmt='0')
box(s,9.0,1.95,3.3,4.35,fill=PANEL,line=LINE)
txt(s,9.25,2.2,2.85,0.4,[[R("INDEX (100 = AVG)",10.5,CYAND,True)]])
txt(s,9.25,2.65,2.85,3.5,[
  [R("Beverages",12,CYAND,True)],
  [R("Peak Jul 136 · trough Feb 74 — widest swing (±62).",10,MUT,False)],
  [R("",7,MUT,False)],
  [R("Frozen Food",12,VIOLET,True)],
  [R("Peaks Jan 122 (winter) — counter-seasonal to drinks.",10,MUT,False)],
  [R("",7,MUT,False)],
  [R("Action: ",11,INK,True),R("phase inventory & promo by category season.",10,MUT,False)],
],sp=1.05)

# =================================================================== SLIDE 9 — FORECAST / SCENARIO
s=slide(); header(s,"FORECAST & WHAT-IF","Reaching €2.5M in 2025 is within reach",9)
txt(s,1.32,1.7,6,0.3,[[R("Monthly revenue + Q1-2025 projection (central +5%)",12,INK,True)]])
m24=D["rev_by_month_24"]
proj=m24+[round(m24[0]*1.05),round(m24[1]*1.05),round(m24[2]*1.05)]
labs=mlab+["Jan'25","Feb'25","Mar'25"]
actser=m24+[None,None,None]
projser=[None]*11+[m24[11]]+[round(m24[0]*1.05),round(m24[1]*1.05),round(m24[2]*1.05)]
add_line(s,1.32,2.1,6.4,4.2,labs,[("Actual 2024",actser),("Projection",projser)],[CYAN,AMBER],fmt='€#,##0,"k"')
box(s,8.1,1.9,4.2,4.4,fill=PANEL,line=LINE)
txt(s,8.35,2.12,3.7,0.4,[[R("SCENARIO MODELLING",10.5,CYAND,True)]])
txt(s,8.35,2.55,3.75,0.95,[[R("Online +20% & Convenience −10%",11,INK,True)],
   [R("net impact on 2024 revenue:",9.5,MUT,False)]],sp=1.0)
txt(s,8.35,3.5,3.75,0.7,[[R("+€68k",34,GREEND,True),R("  +2.9%",15,GREEN,True)]])
box(s,8.35,4.35,3.7,0.02,fill=LINE,radius=False)
txt(s,8.35,4.5,3.75,1.7,[
  [R("To hit €2.5M in 2025",11,INK,True)],
  [R("requires ",10,MUT,False),R("+5.8%",13,CYAND,True),R(" growth — just above the +5.4% delivered in 2024.",10,MUT,False)],
  [R("",5,MUT,False)],
  [R("The Online push alone bridges most of that gap.",10,MUT,False)],
],sp=1.05)

# =================================================================== SLIDE 10 — RECOMMENDATIONS
s=slide(); header(s,"RECOMMENDATIONS","Three actions, quantified",10)
recs=[("01","ACT NOW",GREEN,"Re-baseline the 2025 budget",
       "Rebuild plan volumes on the actual run-rate, not the +5% top-down assumption that missed every cell.",
       "Removes the chronic ≈€115k/yr unfavourable variance"),
      ("02","GROW",CYAN,"Push Supermarket & fix Online",
       "Reinvest behind the fastest channel and rehabilitate Online with promo & assortment.",
       "Online +20% nets +€68k (+2.9%) — covers the gap to €2.5M"),
      ("03","PROTECT",AMBER,"Fix Frozen margin & hedge GBP",
       "Renegotiate Frozen sourcing to the 35% target; hedge sterling on the 27%-of-revenue UK book.",
       "+2pp on €1.42M ≈ +€28k GM; caps a ~6% FX swing")]
for i,(n,tag,col,h,p,imp) in enumerate(recs):
    y=1.75+i*1.55
    box(s,1.32,y,11.0,1.4,fill=WHITE,line=LINE,lw=0.9)
    box(s,1.32,y,0.09,1.4,fill=col,radius=False)
    txt(s,1.6,y+0.2,0.9,1.0,[[R(n,30,RGBColor(0xE3,0xE0,0xDB),True)]])
    chip(s,2.55,y+0.22,tag,col,1.35)
    txt(s,2.55,y+0.66,5.0,0.6,[[R(h,15,INK,True)]])
    txt(s,4.1,y+0.2,4.5,1.05,[[R(p,10.5,MUT,False)]],anchor=MSO_ANCHOR.TOP)
    # impact pill
    box(s,8.7,y+0.3,3.4,0.8,fill=PANEL,line=col,lw=1.0)
    txt(s,8.85,y+0.36,3.1,0.7,[[R("▲ IMPACT  ",9,col,True)],[R(imp,9.5,INK,False)]],sp=1.0)

# =================================================================== SLIDE 11 — RISKS + CLOSE
s=slide(); header(s,"RISKS & CLOSE","What could go wrong — and the bottom line",11)
risks=[("Budget premise",RED,"Assumes the uniform 5% gap is planning optimism, not demand erosion. Monitor the volume trend monthly."),
       ("Scenario independence",AMBER,"The +€68k uplift treats segments as independent; real Online/Convenience cannibalisation would reduce it."),
       ("Margin & FX",VIOLET,"Forecasts hold margin at 38.7% and assume no FX shock. EUR/GBP swung ~6% — a weaker pound compresses 27% of revenue.")]
for i,(a,c,b) in enumerate(risks):
    y=1.8+i*1.05
    box(s,1.32,y,6.0,0.92,fill=PANEL,line=LINE)
    box(s,1.32,y,0.1,0.92,fill=c,radius=False)
    txt(s,1.58,y+0.12,5.6,0.3,[[R(a,11.5,INK,True)]])
    txt(s,1.58,y+0.4,5.6,0.5,[[R(b,9.5,MUT,False)]])
box(s,7.7,1.8,4.6,3.25,fill=None,line=None)
grad(box(s,7.7,1.8,4.6,3.25,fill=CYAN),CYAN,RGBColor(0x22,0xD3,0xEE),45)
txt(s,8.0,2.1,4.0,0.4,[[R("THE BOTTOM LINE",10.5,WHITE,True)]])
txt(s,8.0,2.6,4.05,2.3,[
   [R("A profitable, growing business",13,WHITE,True)],
   [R("held back by an over-optimistic plan, not by execution.",11,WHITE,False)],
   [R("",7,WHITE,False)],
   [R("Rebaseline, back the winners, protect margin → €2.5M in 2025.",11,WHITE,True)],
],sp=1.06)
txt(s,1.32,5.55,11,1.2,[[R("Thank you.  ",22,INK,True),R("Questions?",22,PINK,True)]])
txt(s,1.34,6.35,11,0.4,[[R("EuroFoods S.A. · FP&A Analytics · Live Power BI dashboard + interactive companion available on screen",10.5,MUT,False)]])

prs.save("deliverables/EuroFoods_Presentation.pptx")
print("Saved deck with",len(prs.slides._sldIdLst),"slides")
